from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import os, uuid, zipfile, tempfile, shutil, subprocess, io
from werkzeug.utils import secure_filename

# Core only — heavy libs imported lazily inside each route
from pypdf import PdfReader, PdfWriter
from PIL import Image

app = Flask(__name__)
CORS(app, origins="*", supports_credentials=True)

UPLOAD_FOLDER = tempfile.mkdtemp()
OUTPUT_FOLDER = tempfile.mkdtemp()
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max

ALLOWED = {'pdf', 'doc', 'docx', 'jpg', 'jpeg', 'png', 'webp', 'bmp', 'gif', 'ppt', 'pptx', 'xls', 'xlsx', 'csv'}

def allowed_file(filename, types=None):
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
    return ext in (types or ALLOWED)

def make_output_path(ext):
    return os.path.join(OUTPUT_FOLDER, f"{uuid.uuid4().hex}.{ext}")

def cleanup(*paths):
    for p in paths:
        try:
            if os.path.isfile(p): os.remove(p)
            elif os.path.isdir(p): shutil.rmtree(p)
        except: pass

def add_cors_headers(response):
    response.headers['Access-Control-Allow-Origin'] = '*'
    response.headers['Access-Control-Allow-Methods'] = 'GET, POST, OPTIONS'
    response.headers['Access-Control-Allow-Headers'] = 'Content-Type'
    return response

@app.after_request
def after_request(response):
    return add_cors_headers(response)

@app.before_request
def handle_options():
    if request.method == 'OPTIONS':
        from flask import Response
        r = Response()
        r.headers['Access-Control-Allow-Origin'] = '*'
        r.headers['Access-Control-Allow-Methods'] = 'GET, POST, OPTIONS'
        r.headers['Access-Control-Allow-Headers'] = 'Content-Type'
        return r


# ─── HEALTH CHECK ───
@app.route('/')
def index():
    return jsonify({"status": "EggyPDF API is running!", "tools": 21})


# ─── 1. MERGE PDF ───
@app.route('/api/merge', methods=['POST', 'OPTIONS'])
def merge_pdf():
    files = request.files.getlist('files')
    if len(files) < 2:
        return jsonify({"error": "Please upload at least 2 PDF files."}), 400

    saved = []
    for f in files:
        if not allowed_file(f.filename, {'pdf'}):
            return jsonify({"error": f"{f.filename} is not a valid PDF."}), 400
        path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
        f.save(path)
        saved.append(path)

    out = make_output_path('pdf')
    try:
        writer = PdfWriter()
        for p in saved:
            reader = PdfReader(p)
            for page in reader.pages:
                writer.add_page(page)
        with open(out, 'wb') as fh:
            writer.write(fh)
    except Exception as e:
        cleanup(*saved)
        return jsonify({"error": f"Merge failed: {str(e)}"}), 500

    cleanup(*saved)
    return send_file(out, as_attachment=True, download_name='merged.pdf', mimetype='application/pdf')


# ─── 2. SPLIT PDF ───
@app.route('/api/split', methods=['POST', 'OPTIONS'])
def split_pdf():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400

    split_type = request.form.get('type', 'all')
    page_range = request.form.get('range', '')

    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)

    try:
        reader = PdfReader(saved)
        total = len(reader.pages)
        zip_buf = io.BytesIO()

        with zipfile.ZipFile(zip_buf, 'w', zipfile.ZIP_DEFLATED) as zf:
            if split_type == 'all':
                for i in range(total):
                    writer = PdfWriter()
                    writer.add_page(reader.pages[i])
                    buf = io.BytesIO()
                    writer.write(buf)
                    zf.writestr(f"page_{i+1}.pdf", buf.getvalue())
            else:
                pages = set()
                for part in page_range.split(','):
                    part = part.strip()
                    if '-' in part:
                        a, b = part.split('-')
                        pages.update(range(int(a)-1, int(b)))
                    elif part.isdigit():
                        pages.add(int(part)-1)
                writer = PdfWriter()
                for idx in sorted(pages):
                    if 0 <= idx < total:
                        writer.add_page(reader.pages[idx])
                buf = io.BytesIO()
                writer.write(buf)
                zf.writestr("split_pages.pdf", buf.getvalue())
    except Exception as e:
        cleanup(saved)
        return jsonify({"error": f"Split failed: {str(e)}"}), 500

    cleanup(saved)
    zip_buf.seek(0)
    return send_file(zip_buf, as_attachment=True, download_name='split_pages.zip', mimetype='application/zip')


# ─── 3. COMPRESS PDF ───
@app.route('/api/compress', methods=['POST', 'OPTIONS'])
def compress_pdf():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400

    level = request.form.get('level', 'medium')
    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)
    out = make_output_path('pdf')

    original_size = os.path.getsize(saved)

    # ── Method 1: Ghostscript (best real compression) ──
    gs_settings = {'low': '/ebook', 'medium': '/ebook', 'high': '/screen'}
    quality = gs_settings.get(level, '/ebook')

    gs_ok = False
    for gs_cmd in ['gs', 'ghostscript']:
        try:
            result = subprocess.run([
                gs_cmd, '-sDEVICE=pdfwrite', '-dCompatibilityLevel=1.4',
                f'-dPDFSETTINGS={quality}', '-dNOPAUSE', '-dQUIET', '-dBATCH',
                '-dDetectDuplicateImages=true', '-dCompressFonts=true',
                f'-sOutputFile={out}', saved
            ], capture_output=True, timeout=120)
            if result.returncode == 0 and os.path.exists(out) and os.path.getsize(out) > 0:
                gs_ok = True
                break
        except Exception:
            continue

    # ── Method 2: pikepdf fallback (if Ghostscript unavailable) ──
    if not gs_ok:
        try:
            with pikepdf.open(saved) as pdf:
                pdf.save(out, compress_streams=True,
                         object_stream_mode=pikepdf.ObjectStreamMode.generate,
                         recompress_flate=True)
        except Exception as e:
            cleanup(saved)
            return jsonify({"error": f"Compression failed: {str(e)}"}), 500

    # If output somehow ended up larger, return the original instead
    try:
        if os.path.exists(out) and os.path.getsize(out) >= original_size:
            shutil.copy(saved, out)
    except Exception:
        pass

    cleanup(saved)
    return send_file(out, as_attachment=True, download_name='compressed.pdf', mimetype='application/pdf')


# ─── 4. PDF TO WORD ───
@app.route('/api/pdf-to-word', methods=['POST', 'OPTIONS'])
def pdf_to_word():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400

    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)
    out = make_output_path('docx')

    try:
        from pdf2docx import Converter
        cv = Converter(saved)
        cv.convert(out, start=0, end=None)
        cv.close()
    except Exception as e:
        cleanup(saved)
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500

    cleanup(saved)
    return send_file(out, as_attachment=True, download_name='converted.docx',
                     mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')


# ─── 5. WORD TO PDF ───
@app.route('/api/word-to-pdf', methods=['POST', 'OPTIONS'])
def word_to_pdf():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'doc', 'docx'}):
        return jsonify({"error": "Please upload a valid Word (.doc or .docx) file."}), 400

    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.docx")
    f.save(saved)
    out = make_output_path('pdf')

    try:
        for cmd in ['libreoffice', 'soffice']:
            result = subprocess.run([
                cmd, '--headless', '--convert-to', 'pdf',
                '--outdir', OUTPUT_FOLDER, saved
            ], capture_output=True, timeout=60)
            expected = os.path.join(OUTPUT_FOLDER,
                        os.path.splitext(os.path.basename(saved))[0] + '.pdf')
            if os.path.exists(expected):
                shutil.move(expected, out)
                break
    except Exception as e:
        cleanup(saved)
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500

    cleanup(saved)
    if not os.path.exists(out):
        return jsonify({"error": "Conversion failed. LibreOffice may not be available."}), 500

    return send_file(out, as_attachment=True, download_name='converted.pdf', mimetype='application/pdf')


# ─── 6. JPG TO PDF ───
@app.route('/api/jpg-to-pdf', methods=['POST', 'OPTIONS'])
def jpg_to_pdf():
    files = request.files.getlist('files')
    if not files:
        return jsonify({"error": "Please upload at least one image."}), 400

    images = []
    saved_paths = []
    try:
        for f in files:
            if not allowed_file(f.filename, {'jpg', 'jpeg', 'png'}):
                return jsonify({"error": f"{f.filename} is not a valid image."}), 400
            path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}_{secure_filename(f.filename)}")
            f.save(path)
            saved_paths.append(path)
            img = Image.open(path).convert('RGB')
            images.append(img)

        out = make_output_path('pdf')
        if len(images) == 1:
            images[0].save(out, 'PDF', resolution=100.0)
        else:
            images[0].save(out, 'PDF', resolution=100.0, save_all=True, append_images=images[1:])
    except Exception as e:
        cleanup(*saved_paths)
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500

    cleanup(*saved_paths)
    return send_file(out, as_attachment=True, download_name='images.pdf', mimetype='application/pdf')


# ─── 7. ADD WATERMARK ───
@app.route('/api/watermark', methods=['POST', 'OPTIONS'])
def add_watermark():
    f = request.files.get('file')
    text = request.form.get('text', 'CONFIDENTIAL')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400

    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)
    wm_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}_wm.pdf")

    try:
        from reportlab.pdfgen import canvas as rl_canvas
        from reportlab.lib.pagesizes import A4
        c = rl_canvas.Canvas(wm_path, pagesize=A4)
        c.setFont("Helvetica-Bold", 48)
        c.setFillColorRGB(0.7, 0.7, 0.7, alpha=0.35)
        c.saveState()
        c.translate(A4[0]/2, A4[1]/2)
        c.rotate(45)
        c.drawCentredString(0, 0, text.upper())
        c.restoreState()
        c.save()

        reader = PdfReader(saved)
        wm_reader = PdfReader(wm_path)
        wm_page = wm_reader.pages[0]
        writer = PdfWriter()
        for page in reader.pages:
            page.merge_page(wm_page)
            writer.add_page(page)

        out = make_output_path('pdf')
        with open(out, 'wb') as fh:
            writer.write(fh)
    except Exception as e:
        cleanup(saved, wm_path)
        return jsonify({"error": f"Watermark failed: {str(e)}"}), 500

    cleanup(saved, wm_path)
    return send_file(out, as_attachment=True, download_name='watermarked.pdf', mimetype='application/pdf')


# ─── 8. PROTECT PDF ───
@app.route('/api/protect', methods=['POST', 'OPTIONS'])
def protect_pdf():
    f = request.files.get('file')
    password = request.form.get('password', '')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400
    if not password or len(password) < 4:
        return jsonify({"error": "Password must be at least 4 characters."}), 400

    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)

    try:
        reader = PdfReader(saved)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        writer.encrypt(password)
        out = make_output_path('pdf')
        with open(out, 'wb') as fh:
            writer.write(fh)
    except Exception as e:
        cleanup(saved)
        return jsonify({"error": f"Protection failed: {str(e)}"}), 500

    cleanup(saved)
    return send_file(out, as_attachment=True, download_name='protected.pdf', mimetype='application/pdf')


# ─── 9. PDF TO JPG ───
@app.route('/api/pdf-to-jpg', methods=['POST', 'OPTIONS'])
def pdf_to_jpg():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400

    dpi = int(request.form.get('dpi', 150))
    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)

    try:
        # Use ghostscript to render pages to images
        out_dir = os.path.join(OUTPUT_FOLDER, uuid.uuid4().hex)
        os.makedirs(out_dir)

        result = subprocess.run([
            'gs', '-dNOPAUSE', '-dBATCH', '-dSAFER',
            '-sDEVICE=jpeg', f'-r{dpi}',
            f'-sOutputFile={out_dir}/page_%03d.jpg',
            saved
        ], capture_output=True, timeout=120)

        jpg_files = sorted([
            os.path.join(out_dir, fn)
            for fn in os.listdir(out_dir)
            if fn.endswith('.jpg')
        ])

        if not jpg_files:
            raise Exception("No pages rendered")

        if len(jpg_files) == 1:
            # Single page — return JPG directly
            cleanup(saved)
            return send_file(jpg_files[0], as_attachment=True,
                           download_name='page_1.jpg', mimetype='image/jpeg')
        else:
            # Multiple pages — zip them
            zip_buf = io.BytesIO()
            with zipfile.ZipFile(zip_buf, 'w', zipfile.ZIP_DEFLATED) as zf:
                for jpg in jpg_files:
                    zf.write(jpg, os.path.basename(jpg))
            zip_buf.seek(0)
            cleanup(saved, out_dir)
            return send_file(zip_buf, as_attachment=True,
                           download_name='pdf_pages.zip', mimetype='application/zip')
    except Exception as e:
        cleanup(saved)
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500


# ─── 10. PDF TO PNG ───
@app.route('/api/pdf-to-png', methods=['POST', 'OPTIONS'])
def pdf_to_png():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400

    dpi = int(request.form.get('dpi', 150))
    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)

    try:
        out_dir = os.path.join(OUTPUT_FOLDER, uuid.uuid4().hex)
        os.makedirs(out_dir)

        subprocess.run([
            'gs', '-dNOPAUSE', '-dBATCH', '-dSAFER',
            '-sDEVICE=png16m', f'-r{dpi}',
            f'-sOutputFile={out_dir}/page_%03d.png',
            saved
        ], capture_output=True, timeout=120)

        png_files = sorted([
            os.path.join(out_dir, fn)
            for fn in os.listdir(out_dir)
            if fn.endswith('.png')
        ])

        if not png_files:
            raise Exception("No pages rendered")

        if len(png_files) == 1:
            cleanup(saved)
            return send_file(png_files[0], as_attachment=True,
                           download_name='page_1.png', mimetype='image/png')
        else:
            zip_buf = io.BytesIO()
            with zipfile.ZipFile(zip_buf, 'w', zipfile.ZIP_DEFLATED) as zf:
                for png in png_files:
                    zf.write(png, os.path.basename(png))
            zip_buf.seek(0)
            cleanup(saved, out_dir)
            return send_file(zip_buf, as_attachment=True,
                           download_name='pdf_pages.zip', mimetype='application/zip')
    except Exception as e:
        cleanup(saved)
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500


# ─── 11. PNG TO PDF ───
@app.route('/api/png-to-pdf', methods=['POST', 'OPTIONS'])
def png_to_pdf():
    files = request.files.getlist('files')
    if not files:
        return jsonify({"error": "Please upload at least one PNG or image file."}), 400

    images = []
    saved_paths = []
    try:
        for f in files:
            if not allowed_file(f.filename, {'jpg', 'jpeg', 'png', 'webp', 'bmp', 'gif'}):
                return jsonify({"error": f"{f.filename} is not a supported image."}), 400
            path = os.path.join(UPLOAD_FOLDER,
                               f"{uuid.uuid4().hex}_{secure_filename(f.filename)}")
            f.save(path)
            saved_paths.append(path)
            img = Image.open(path).convert('RGB')
            images.append(img)

        out = make_output_path('pdf')
        if len(images) == 1:
            images[0].save(out, 'PDF', resolution=150.0)
        else:
            images[0].save(out, 'PDF', resolution=150.0,
                          save_all=True, append_images=images[1:])
    except Exception as e:
        cleanup(*saved_paths)
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500

    cleanup(*saved_paths)
    return send_file(out, as_attachment=True,
                    download_name='converted.pdf', mimetype='application/pdf')


# ─── 12. ROTATE PDF ───
@app.route('/api/rotate', methods=['POST', 'OPTIONS'])
def rotate_pdf():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400

    angle = int(request.form.get('angle', 90))
    if angle not in [90, 180, 270]:
        return jsonify({"error": "Angle must be 90, 180, or 270."}), 400

    pages_input = request.form.get('pages', 'all')

    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)

    try:
        reader = PdfReader(saved)
        writer = PdfWriter()
        total = len(reader.pages)

        # Parse which pages to rotate
        if pages_input == 'all':
            rotate_pages = set(range(total))
        else:
            rotate_pages = set()
            for part in pages_input.split(','):
                part = part.strip()
                if '-' in part:
                    a, b = part.split('-')
                    rotate_pages.update(range(int(a)-1, int(b)))
                elif part.isdigit():
                    rotate_pages.add(int(part)-1)

        for i, page in enumerate(reader.pages):
            if i in rotate_pages:
                page.rotate(angle)
            writer.add_page(page)

        out = make_output_path('pdf')
        with open(out, 'wb') as fh:
            writer.write(fh)
    except Exception as e:
        cleanup(saved)
        return jsonify({"error": f"Rotation failed: {str(e)}"}), 500

    cleanup(saved)
    return send_file(out, as_attachment=True,
                    download_name='rotated.pdf', mimetype='application/pdf')


# ─── 13. DELETE PDF PAGES ───
@app.route('/api/delete-pages', methods=['POST', 'OPTIONS'])
def delete_pages():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400

    pages_input = request.form.get('pages', '')
    if not pages_input:
        return jsonify({"error": "Please specify pages to delete."}), 400

    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)

    try:
        reader = PdfReader(saved)
        total = len(reader.pages)

        # Parse pages to delete
        delete_set = set()
        for part in pages_input.split(','):
            part = part.strip()
            if '-' in part:
                a, b = part.split('-')
                delete_set.update(range(int(a)-1, int(b)))
            elif part.isdigit():
                delete_set.add(int(part)-1)

        writer = PdfWriter()
        kept = 0
        for i, page in enumerate(reader.pages):
            if i not in delete_set:
                writer.add_page(page)
                kept += 1

        if kept == 0:
            cleanup(saved)
            return jsonify({"error": "Cannot delete all pages from a PDF."}), 400

        out = make_output_path('pdf')
        with open(out, 'wb') as fh:
            writer.write(fh)
    except Exception as e:
        cleanup(saved)
        return jsonify({"error": f"Page deletion failed: {str(e)}"}), 500

    cleanup(saved)
    return send_file(out, as_attachment=True,
                    download_name='edited.pdf', mimetype='application/pdf')


# ─── 14. UNLOCK PDF ───
@app.route('/api/unlock', methods=['POST', 'OPTIONS'])
def unlock_pdf():
    f = request.files.get('file')
    password = request.form.get('password', '')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400

    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)
    out = make_output_path('pdf')

    opened = False

    # ── Strategy 1: pikepdf (handles most encrypted PDFs) ──
    try:
        import pikepdf
        pwds = [password] if password else []
        pwds += ['', ' ']
        for pwd in pwds:
            try:
                pdf = pikepdf.open(saved, password=pwd)
                # Save without encryption. allow_overwriting_input not needed (different path)
                pdf.save(out)
                pdf.close()
                opened = True
                break
            except pikepdf.PasswordError:
                continue
            except Exception:
                # Try saving with normalized content if direct save fails
                try:
                    pdf = pikepdf.open(saved, password=pwd)
                    pdf.save(out, fix_metadata_version=True)
                    pdf.close()
                    opened = True
                    break
                except Exception:
                    continue
    except Exception:
        pass

    # ── Strategy 2: pypdf fallback ──
    if not opened:
        try:
            reader = PdfReader(saved)
            if reader.is_encrypted:
                # decrypt returns 0=fail, 1=user pw, 2=owner pw
                result = 0
                for pwd in ([password] if password else []) + ['', ' ']:
                    try:
                        result = reader.decrypt(pwd)
                        if result != 0:
                            break
                    except Exception:
                        continue
                if result == 0:
                    cleanup(saved)
                    if password:
                        return jsonify({"error": "Incorrect password. Please double-check and try again. If you have forgotten the password, it cannot be recovered."}), 400
                    return jsonify({"error": "This PDF requires a password to open. Please enter the correct password. Note: if you do not know the password, it cannot be removed — this is a security feature of PDF encryption."}), 400

            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            with open(out, 'wb') as fh:
                writer.write(fh)
            opened = True
        except Exception as e:
            cleanup(saved)
            return jsonify({"error": f"Unlock failed: {str(e)}"}), 500

    if not opened:
        cleanup(saved)
        if password:
            return jsonify({"error": "Incorrect password. Please double-check and try again. If you have forgotten the password, it cannot be recovered."}), 400
        return jsonify({"error": "This PDF requires a password to open. Please enter the correct password. Note: if you do not know the password, it cannot be removed — this is a security feature of PDF encryption."}), 400

    if not os.path.exists(out) or os.path.getsize(out) == 0:
        cleanup(saved)
        return jsonify({"error": "Unlock failed. The file may be corrupted."}), 500

    cleanup(saved)
    return send_file(out, as_attachment=True,
                    download_name='unlocked.pdf', mimetype='application/pdf')



# ─── 15. ADD PAGE NUMBERS ───
@app.route('/api/page-numbers', methods=['POST', 'OPTIONS'])
def add_page_numbers():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400

    position   = request.form.get('position', 'bottom-center')
    start_from = int(request.form.get('start_from', 1))
    font_size  = int(request.form.get('font_size', 12))
    prefix     = request.form.get('prefix', '')

    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)

    try:
        from reportlab.pdfgen import canvas as rl_canvas
        from reportlab.lib.pagesizes import A4

        reader  = PdfReader(saved)
        writer  = PdfWriter()
        total   = len(reader.pages)

        for i, page in enumerate(reader.pages):
            w = float(page.mediabox.width)
            h = float(page.mediabox.height)

            # Build a single-page overlay with the page number
            overlay_buf = io.BytesIO()
            c = rl_canvas.Canvas(overlay_buf, pagesize=(w, h))
            c.setFont("Helvetica", font_size)
            c.setFillColorRGB(0.3, 0.3, 0.3)

            label = f"{prefix}{i + start_from}"

            margin = 28
            pos_map = {
                'bottom-center': (w / 2,       margin),
                'bottom-left':   (margin,       margin),
                'bottom-right':  (w - margin,   margin),
                'top-center':    (w / 2,       h - margin),
                'top-left':      (margin,       h - margin),
                'top-right':     (w - margin,  h - margin),
            }
            x, y = pos_map.get(position, (w / 2, margin))

            if 'center' in position:
                c.drawCentredString(x, y, label)
            elif 'right' in position:
                c.drawRightString(x, y, label)
            else:
                c.drawString(x, y, label)

            c.save()
            overlay_buf.seek(0)

            overlay_page = PdfReader(overlay_buf).pages[0]
            page.merge_page(overlay_page)
            writer.add_page(page)

        out = make_output_path('pdf')
        with open(out, 'wb') as fh:
            writer.write(fh)

    except Exception as e:
        cleanup(saved)
        return jsonify({"error": f"Failed to add page numbers: {str(e)}"}), 500

    cleanup(saved)
    return send_file(out, as_attachment=True,
                     download_name='numbered.pdf', mimetype='application/pdf')


# ─── 16. CROP PDF ───
@app.route('/api/crop', methods=['POST', 'OPTIONS'])
def crop_pdf():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400

    try:
        top    = float(request.form.get('top',    0))
        bottom = float(request.form.get('bottom', 0))
        left   = float(request.form.get('left',   0))
        right  = float(request.form.get('right',  0))
    except ValueError:
        return jsonify({"error": "Crop values must be numbers."}), 400

    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)

    try:
        reader = PdfReader(saved)
        writer = PdfWriter()

        for page in reader.pages:
            mb = page.mediabox
            new_x0 = float(mb.lower_left[0])  + left
            new_y0 = float(mb.lower_left[1])  + bottom
            new_x1 = float(mb.upper_right[0]) - right
            new_y1 = float(mb.upper_right[1]) - top

            if new_x1 <= new_x0 or new_y1 <= new_y0:
                cleanup(saved)
                return jsonify({"error": "Crop margins are too large for this page size."}), 400

            page.mediabox.lower_left  = (new_x0, new_y0)
            page.mediabox.upper_right = (new_x1, new_y1)
            writer.add_page(page)

        out = make_output_path('pdf')
        with open(out, 'wb') as fh:
            writer.write(fh)

    except Exception as e:
        cleanup(saved)
        return jsonify({"error": f"Crop failed: {str(e)}"}), 500

    cleanup(saved)
    return send_file(out, as_attachment=True,
                     download_name='cropped.pdf', mimetype='application/pdf')


# ─── 17. POWERPOINT TO PDF ───
@app.route('/api/ppt-to-pdf', methods=['POST', 'OPTIONS'])
def ppt_to_pdf():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'ppt', 'pptx'}):
        return jsonify({"error": "Please upload a valid PowerPoint (.ppt or .pptx) file."}), 400

    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pptx")
    f.save(saved)
    out = make_output_path('pdf')

    try:
        for cmd in ['libreoffice', 'soffice']:
            result = subprocess.run([
                cmd, '--headless', '--convert-to', 'pdf',
                '--outdir', OUTPUT_FOLDER, saved
            ], capture_output=True, timeout=120)
            expected = os.path.join(
                OUTPUT_FOLDER,
                os.path.splitext(os.path.basename(saved))[0] + '.pdf'
            )
            if os.path.exists(expected):
                shutil.move(expected, out)
                break
    except Exception as e:
        cleanup(saved)
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500

    cleanup(saved)
    if not os.path.exists(out):
        return jsonify({"error": "Conversion failed. Please try again."}), 500

    return send_file(out, as_attachment=True,
                     download_name='converted.pdf', mimetype='application/pdf')


# ─── 18. EXCEL TO PDF ───
@app.route('/api/excel-to-pdf', methods=['POST', 'OPTIONS'])
def excel_to_pdf():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'xls', 'xlsx', 'csv'}):
        return jsonify({"error": "Please upload a valid Excel (.xls, .xlsx) or CSV file."}), 400

    ext   = f.filename.rsplit('.', 1)[-1].lower()
    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.{ext}")
    f.save(saved)
    out   = make_output_path('pdf')

    try:
        for cmd in ['libreoffice', 'soffice']:
            result = subprocess.run([
                cmd, '--headless', '--convert-to', 'pdf',
                '--outdir', OUTPUT_FOLDER, saved
            ], capture_output=True, timeout=120)
            expected = os.path.join(
                OUTPUT_FOLDER,
                os.path.splitext(os.path.basename(saved))[0] + '.pdf'
            )
            if os.path.exists(expected):
                shutil.move(expected, out)
                break
    except Exception as e:
        cleanup(saved)
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500

    cleanup(saved)
    if not os.path.exists(out):
        return jsonify({"error": "Conversion failed. Please try again."}), 500

    return send_file(out, as_attachment=True,
                     download_name='converted.pdf', mimetype='application/pdf')


# ─── 19. PDF TO EXCEL ───
@app.route('/api/pdf-to-excel', methods=['POST', 'OPTIONS'])
def pdf_to_excel():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400

    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)
    out   = make_output_path('xlsx')

    try:
        import pdfplumber
        import openpyxl

        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # remove default sheet
        found_any = False

        with pdfplumber.open(saved) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables()
                if tables:
                    for t_idx, table in enumerate(tables):
                        ws = wb.create_sheet(title=f"Page{page_num}_T{t_idx+1}")
                        for row in table:
                            ws.append([cell or '' for cell in row])
                        found_any = True
                else:
                    # No table — extract raw text into a sheet
                    text = page.extract_text() or ''
                    if text.strip():
                        ws = wb.create_sheet(title=f"Page{page_num}_Text")
                        for line in text.split('\n'):
                            ws.append([line])
                        found_any = True

        if not found_any:
            cleanup(saved)
            return jsonify({"error": "No extractable content found in this PDF."}), 400

        wb.save(out)

    except ImportError:
        cleanup(saved)
        return jsonify({"error": "PDF to Excel library not installed on server."}), 500
    except Exception as e:
        cleanup(saved)
        return jsonify({"error": f"Extraction failed: {str(e)}"}), 500

    cleanup(saved)
    return send_file(
        out, as_attachment=True,
        download_name='extracted.xlsx',
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )


# ─── 20. PDF TO TEXT ───
@app.route('/api/pdf-to-text', methods=['POST', 'OPTIONS'])
def pdf_to_text():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400

    saved = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)

    try:
        reader = PdfReader(saved)
        lines  = []
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ''
            if text.strip():
                lines.append(f"--- Page {i} ---\n{text.strip()}")

        if not lines:
            cleanup(saved)
            return jsonify({"error": "No readable text found in this PDF. It may be a scanned image."}), 400

        full_text = '\n\n'.join(lines)
        out_buf   = io.BytesIO(full_text.encode('utf-8'))
        out_buf.seek(0)

    except Exception as e:
        cleanup(saved)
        return jsonify({"error": f"Text extraction failed: {str(e)}"}), 500

    cleanup(saved)
    return send_file(out_buf, as_attachment=True,
                     download_name='extracted.txt', mimetype='text/plain')


# ─── 21. PDF TO POWERPOINT ───
@app.route('/api/pdf-to-ppt', methods=['POST', 'OPTIONS'])
def pdf_to_ppt():
    f = request.files.get('file')
    if not f or not allowed_file(f.filename, {'pdf'}):
        return jsonify({"error": "Please upload a valid PDF file."}), 400

    saved   = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex}.pdf")
    f.save(saved)
    out_dir = os.path.join(OUTPUT_FOLDER, uuid.uuid4().hex)
    os.makedirs(out_dir)
    out     = make_output_path('pptx')

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        import re

        # Render each page as image via ghostscript
        subprocess.run([
            'gs', '-dNOPAUSE', '-dBATCH', '-dSAFER',
            '-sDEVICE=png16m', '-r150',
            f'-sOutputFile={out_dir}/slide_%03d.png', saved
        ], capture_output=True, timeout=180)

        slides = sorted([
            os.path.join(out_dir, fn)
            for fn in os.listdir(out_dir) if fn.endswith('.png')
        ])

        if not slides:
            raise Exception("Could not render PDF pages")

        prs = Presentation()
        prs.slide_width  = Inches(10)
        prs.slide_height = Inches(7.5)
        blank_layout     = prs.slide_layouts[6]  # blank

        for slide_img in slides:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                slide_img, Inches(0), Inches(0),
                width=Inches(10), height=Inches(7.5)
            )

        prs.save(out)

    except ImportError:
        cleanup(saved, out_dir)
        return jsonify({"error": "python-pptx not installed on server."}), 500
    except Exception as e:
        cleanup(saved, out_dir)
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500

    cleanup(saved, out_dir)
    return send_file(
        out, as_attachment=True,
        download_name='converted.pptx',
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )




# ─── SEND CV TO EMAIL ───
@app.route('/api/send-cv-email', methods=['POST', 'OPTIONS'])
def send_cv_email():
    """
    Sends CV email with a perfectly formatted PDF attachment.
    The PDF is generated from the exact same HTML+CSS the browser uses,
    ensuring the email PDF looks identical to the downloaded version.
    Uses Pillow to crop the photo into a circle before embedding in PDF.
    """
    import urllib.request, urllib.error, json as json_lib, base64, io

    BREVO_API_KEY    = os.environ.get('BREVO_API_KEY', '').strip()
    SENDER_EMAIL     = os.environ.get('BREVO_SENDER_EMAIL', 'hello@eggypdf.com').strip()
    if not BREVO_API_KEY:
        return jsonify({"error": "Email service not configured."}), 503

    data         = request.get_json(silent=True) or {}
    email        = (data.get('email') or '').strip()
    name         = (data.get('name') or 'there').strip()
    resume_html  = (data.get('resume_html') or '').strip()
    template_css = (data.get('template_css') or '').strip()
    photo_data   = (data.get('photo_data') or '').strip()
    photo_pos_x  = float(data.get('photo_pos_x') or 50)
    photo_pos_y  = float(data.get('photo_pos_y') or 50)

    print(f"Email request: name={name}, email={email}, photo={'YES' if photo_data else 'NO'}")

    if not email or not resume_html:
        return jsonify({"error": "Email and resume content required."}), 400

    safe_name = (name.replace(' ', '_') or 'Resume')
    pdf_b64 = None

    # ── Generate PDF using reportlab ──
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.units import cm, mm
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable, Table, TableStyle, Image as RLImage
        import re as _re

        def st(h):
            """Strip HTML tags and decode entities"""
            t = _re.sub(r'<[^>]+>', ' ', str(h))
            for e, r in [('&amp;','&'),('&lt;','<'),('&gt;','>'),('&nbsp;',' '),
                         ('&#39;',"'"),('&quot;','"'),('&apos;',"'"),('–','-'),('—','-')]:
                t = t.replace(e, r)
            return ' '.join(t.split()).strip()

        def xs(t):
            """Escape for reportlab XML"""
            return str(t).replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')

        def fnd(pat, html):
            m = _re.search(pat, html, _re.DOTALL|_re.IGNORECASE)
            return st(m.group(1)) if m else ''

        def fall(pat, html):
            return _re.findall(pat, html, _re.DOTALL|_re.IGNORECASE)

        # ── Styles matching the browser CSS ──
        DARK  = colors.HexColor('#1a1a2e')
        DGREY = colors.HexColor('#4a5568')
        GREY  = colors.HexColor('#6b7280')
        LGREY = colors.HexColor('#e5e7eb')
        BODY  = colors.HexColor('#374151')

        sName  = ParagraphStyle('N',  fontName='Helvetica-Bold', fontSize=22, textColor=DARK,  spaceAfter=4,  leading=26)
        sJob   = ParagraphStyle('J',  fontName='Helvetica',      fontSize=11, textColor=DGREY, spaceAfter=0,  leading=16)
        sHead  = ParagraphStyle('H',  fontName='Helvetica-Bold', fontSize=9,  textColor=DARK,  spaceBefore=14,spaceAfter=5, leading=12)
        sBody  = ParagraphStyle('B',  fontName='Helvetica',      fontSize=9,  textColor=BODY,  leading=14,   spaceAfter=3)
        sBold  = ParagraphStyle('BB', fontName='Helvetica-Bold', fontSize=9,  textColor=DARK,  leading=14,   spaceAfter=2)
        sSub   = ParagraphStyle('S',  fontName='Helvetica',      fontSize=9,  textColor=GREY,  leading=13,   spaceAfter=2)
        sDate  = ParagraphStyle('D',  fontName='Helvetica',      fontSize=8,  textColor=GREY,  leading=12,   spaceAfter=2)
        sBul   = ParagraphStyle('BU', fontName='Helvetica',      fontSize=9,  textColor=BODY,  leading=14,   leftIndent=10, spaceAfter=2)
        sSkill = ParagraphStyle('SK', fontName='Helvetica',      fontSize=9,  textColor=BODY,  leading=14,   spaceAfter=3)

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4,
                                leftMargin=1.8*cm, rightMargin=1.8*cm,
                                topMargin=1.8*cm, bottomMargin=1.8*cm)
        story = []

        # ── NAME + PHOTO ──
        cv_name = fnd(r'class="[^"]*resume-name[^"]*"[^>]*>(.*?)</', resume_html) or name
        cv_job  = fnd(r'class="[^"]*resume-job-title[^"]*"[^>]*>(.*?)</', resume_html)
        if not cv_job:
            cv_job = fnd(r'class="[^"]*resume-title[^"]*"[^>]*>(.*?)</', resume_html)

        if photo_data and photo_data.startswith('data:image'):
            try:
                from PIL import Image as _PIL, ImageDraw as _Draw
                _, enc = photo_data.split(',', 1)
                raw = base64.b64decode(enc)
                pil = _PIL.open(io.BytesIO(raw)).convert('RGBA')
                # High quality: work at 6x then scale down
                target_pt = 90
                work_sz = 540
                pil = pil.resize((work_sz, work_sz), _PIL.LANCZOS)
                # Circle mask
                mask = _PIL.new('L', (work_sz, work_sz), 0)
                _Draw.Draw(mask).ellipse((0, 0, work_sz, work_sz), fill=255)
                circle = _PIL.new('RGBA', (work_sz, work_sz), (255,255,255,0))
                circle.paste(pil, mask=mask)
                # Draw dark navy border
                border_draw = _Draw.Draw(circle)
                bw = 10
                border_draw.ellipse((bw, bw, work_sz-bw, work_sz-bw),
                    outline=(26, 26, 46, 255), width=bw)
                # Flatten to white background
                final = _PIL.new('RGB', (work_sz, work_sz), (255,255,255))
                final.paste(circle, mask=circle.split()[3])
                final = final.resize((target_pt*3, target_pt*3), _PIL.LANCZOS)
                cb = io.BytesIO()
                final.save(cb, format='PNG')
                cb.seek(0)
                photo_img = RLImage(cb, width=target_pt, height=target_pt)

                name_col = [Paragraph(xs(cv_name), sName)]
                if cv_job: name_col.append(Paragraph(xs(cv_job), sJob))
                tbl = Table([[photo_img, name_col]], colWidths=[80, None])
                tbl.setStyle(TableStyle([
                    ('VALIGN',(0,0),(-1,-1),'MIDDLE'),
                    ('LEFTPADDING',(0,0),(-1,-1),0),
                    ('RIGHTPADDING',(0,0),(-1,-1),8),
                    ('TOPPADDING',(0,0),(-1,-1),0),
                    ('BOTTOMPADDING',(0,0),(-1,-1),6),
                ]))
                story.append(tbl)
                print("Photo added to PDF in circle")
            except Exception as pe:
                print(f"Photo error: {pe}")
                story.append(Paragraph(xs(cv_name), sName))
                if cv_job: story.append(Paragraph(xs(cv_job), sJob))
        else:
            story.append(Paragraph(xs(cv_name), sName))
            if cv_job: story.append(Paragraph(xs(cv_job), sJob))

        story.append(HRFlowable(width='100%', thickness=2, color=DARK, spaceAfter=10))

        # ── CONTACT ──
        contact_rows = fall(r'<div class="[^"]*contact-row[^"]*">(.*?)</div>', resume_html)
        if contact_rows:
            story.append(Paragraph('CONTACT', sHead))
            story.append(HRFlowable(width='100%', thickness=0.5, color=LGREY, spaceAfter=5))
            # Two-column contact layout matching browser
            rows = []
            for row in contact_rows:
                lbl = fnd(r'class="[^"]*contact-lbl[^"]*"[^>]*>(.*?)</span>', row)
                val = fnd(r'class="[^"]*contact-val[^"]*"[^>]*>(.*?)</span>', row)
                if lbl and val:
                    rows.append((lbl, val))
            # Display in 2 columns
            for i in range(0, len(rows), 2):
                left  = rows[i]
                right = rows[i+1] if i+1 < len(rows) else None
                left_cell  = Paragraph(f'<b>{xs(left[0])}:</b>  {xs(left[1])}', sBody)
                right_cell = Paragraph(f'<b>{xs(right[0])}:</b>  {xs(right[1])}', sBody) if right else Paragraph('', sBody)
                tbl = Table([[left_cell, right_cell]], colWidths=['50%','50%'])
                tbl.setStyle(TableStyle([
                    ('VALIGN',(0,0),(-1,-1),'TOP'),
                    ('LEFTPADDING',(0,0),(-1,-1),0),
                    ('RIGHTPADDING',(0,0),(-1,-1),0),
                    ('TOPPADDING',(0,0),(-1,-1),2),
                    ('BOTTOMPADDING',(0,0),(-1,-1),2),
                ]))
                story.append(tbl)

        # ── ALL OTHER SECTIONS ──
        headings = fall(r'<div class="[^"]*section-heading[^"]*">(.*?)</div>', resume_html)
        parts    = _re.split(r'<div class="[^"]*section-heading[^"]*">.*?</div>', resume_html, flags=_re.DOTALL|_re.IGNORECASE)

        for i, heading in enumerate(headings):
            htxt = st(heading).strip()
            HTXT = htxt.upper()
            block = parts[i+1] if i+1 < len(parts) else ''

            if 'CONTACT' in HTXT:
                continue  # already handled above

            story.append(Paragraph(xs(HTXT), sHead))
            story.append(HRFlowable(width='100%', thickness=0.5, color=LGREY, spaceAfter=5))

            if any(k in HTXT for k in ('SUMMARY','PROFILE','ABOUT')):
                txt = fnd(r'class="[^"]*summary-text[^"]*"[^>]*>(.*?)</p>', block)
                if not txt: txt = st(block[:800])
                if txt: story.append(Paragraph(xs(txt), sBody))

            elif any(k in HTXT for k in ('SKILL','LANGUAGE')):
                pills = fall(r'<span class="[^"]*skill-pill[^"]*">(.*?)</span>', block)
                if pills:
                    from reportlab.platypus import Table, TableStyle
                    from reportlab.lib import colors as _colors

                    sPill = ParagraphStyle('Pill',
                        fontName='Helvetica', fontSize=8,
                        textColor=DARK, leading=12)

                    # Build pill cells
                    pill_cells = []
                    for p in pills:
                        txt = xs(st(p))
                        cell = Paragraph(txt, sPill)
                        pill_cells.append(cell)

                    # Arrange pills in rows of up to 5
                    row_size = 5
                    rows = [pill_cells[i:i+row_size] for i in range(0, len(pill_cells), row_size)]

                    for row in rows:
                        # Pad row to row_size
                        while len(row) < row_size:
                            row.append(Paragraph('', sPill))
                        col_w = (doc.width) / row_size
                        tbl = Table([row], colWidths=[col_w]*row_size)
                        tbl.setStyle(TableStyle([
                            ('BOX',       (0,0), (0,-1), 0.5, _colors.HexColor('#1a1a2e')),
                            ('BOX',       (1,0), (1,-1), 0.5, _colors.HexColor('#1a1a2e')),
                            ('BOX',       (2,0), (2,-1), 0.5, _colors.HexColor('#1a1a2e')),
                            ('BOX',       (3,0), (3,-1), 0.5, _colors.HexColor('#1a1a2e')),
                            ('BOX',       (4,0), (4,-1), 0.5, _colors.HexColor('#1a1a2e')),
                            ('LEFTPADDING',  (0,0), (-1,-1), 6),
                            ('RIGHTPADDING', (0,0), (-1,-1), 6),
                            ('TOPPADDING',   (0,0), (-1,-1), 4),
                            ('BOTTOMPADDING',(0,0), (-1,-1), 4),
                            ('VALIGN',       (0,0), (-1,-1), 'MIDDLE'),
                            ('ALIGN',        (0,0), (-1,-1), 'CENTER'),
                        ]))
                        story.append(tbl)
                        story.append(Spacer(1, 4))

            elif any(k in HTXT for k in ('EXPERIENCE','WORK')):
                entries = _re.findall(
                    r'<div class="[^"]*resume-entry[^"]*">(.*?)</div>\s*</div>',
                    block, _re.DOTALL
                )
                for entry in entries:
                    et    = fnd(r'class="entry-title"[^>]*>(.*?)</div>', entry)
                    spans = fall(r'<span>(.*?)</span>', entry)
                    buls  = fall(r'<li[^>]*>(.*?)</li>', entry)
                    desc  = fnd(r'class="[^"]*entry-desc[^"]*"[^>]*>(.*?)</p>', entry)

                    if et:
                        story.append(Paragraph(f'<b>{xs(et)}</b>', sBold))
                    if len(spans) >= 2:
                        company = xs(st(spans[0]))
                        date    = xs(st(spans[1]))
                        # Company left, date right
                        tbl = Table(
                            [[Paragraph(company, sSub), Paragraph(date, sDate)]],
                            colWidths=['70%','30%']
                        )
                        tbl.setStyle(TableStyle([
                            ('VALIGN',(0,0),(-1,-1),'TOP'),
                            ('LEFTPADDING',(0,0),(-1,-1),0),
                            ('RIGHTPADDING',(0,0),(-1,-1),0),
                            ('TOPPADDING',(0,0),(-1,-1),0),
                            ('BOTTOMPADDING',(0,0),(-1,-1),3),
                            ('ALIGN',(1,0),(1,0),'RIGHT'),
                        ]))
                        story.append(tbl)
                    elif len(spans) == 1:
                        story.append(Paragraph(xs(st(spans[0])), sSub))

                    for bl in buls:
                        story.append(Paragraph(f'• {xs(st(bl))}', sBul))
                    if desc and not buls:
                        story.append(Paragraph(xs(st(desc)), sBul))
                    story.append(Spacer(1, 5))

            elif 'EDUCATION' in HTXT:
                entries = _re.findall(
                    r'<div class="[^"]*resume-entry[^"]*">(.*?)</div>\s*</div>',
                    block, _re.DOTALL
                )
                for entry in entries:
                    et    = fnd(r'class="entry-title"[^>]*>(.*?)</div>', entry)
                    spans = fall(r'<span>(.*?)</span>', entry)
                    if et: story.append(Paragraph(f'<b>{xs(et)}</b>', sBold))
                    if len(spans) >= 2:
                        tbl = Table(
                            [[Paragraph(xs(st(spans[0])), sSub), Paragraph(xs(st(spans[1])), sDate)]],
                            colWidths=['70%','30%']
                        )
                        tbl.setStyle(TableStyle([
                            ('VALIGN',(0,0),(-1,-1),'TOP'),
                            ('LEFTPADDING',(0,0),(-1,-1),0),
                            ('RIGHTPADDING',(0,0),(-1,-1),0),
                            ('TOPPADDING',(0,0),(-1,-1),0),
                            ('BOTTOMPADDING',(0,0),(-1,-1),3),
                            ('ALIGN',(1,0),(1,0),'RIGHT'),
                        ]))
                        story.append(tbl)
                    elif len(spans) == 1:
                        story.append(Paragraph(xs(st(spans[0])), sSub))
                    story.append(Spacer(1, 4))

            else:
                txt = st(block[:400])
                if txt: story.append(Paragraph(xs(txt), sBody))

        doc.build(story)
        pdf_b64 = base64.b64encode(buf.getvalue()).decode('utf-8')
        print(f"PDF generated OK: {len(buf.getvalue())} bytes")

    except Exception as e:
        import traceback
        pdf_b64 = None
        print(f"PDF error: {e}")
        print(traceback.format_exc())

    # ── Email body ──
    email_body = f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8"/></head>
<body style="margin:0;padding:0;background:#f3f4f6;font-family:Arial,sans-serif">
  <div style="background:#1a1a2e;padding:22px 24px;text-align:center">
    <div style="font-size:1.3rem;font-weight:700;color:#fff">Eggy<span style="color:#f5a623">PDF</span></div>
    <p style="color:rgba(255,255,255,0.7);font-size:13px;margin:6px 0 0">Your CV is ready, {name}!</p>
  </div>
  <div style="max-width:600px;margin:0 auto;padding:24px 16px">
    <div style="background:#fff8ed;border:2px solid #f5a623;border-radius:16px;padding:20px 24px;margin-bottom:20px;text-align:center">
      <div style="font-size:2rem;margin-bottom:8px">🎉</div>
      <p style="font-size:14px;color:#92400e;line-height:1.7;margin:0">
        {"📎 Your <strong>CV is attached as a PDF</strong> — open it, save it, share it with employers." if pdf_b64 else "Please download your CV directly from EggyPDF."}
      </p>
    </div>
    <div style="background:#fff;border-radius:12px;padding:20px 24px;margin-bottom:20px;border:1px solid #e5e7eb">
      <p style="font-size:13px;color:#374151;line-height:1.8;margin:0">
        ✅ You are on our <strong>early access list</strong> for upcoming AI features.<br/>
        We will notify you first when new tools go live — completely free.
      </p>
    </div>
    <div style="text-align:center;margin-bottom:24px">
      <a href="https://eggypdf.com/resume-builder.html"
         style="display:inline-block;background:#f5a623;color:#fff;padding:13px 32px;border-radius:50px;text-decoration:none;font-weight:700;font-size:14px">
        ✏️ Edit my CV on EggyPDF
      </a>
    </div>
    <p style="font-size:11px;color:#9ca3af;text-align:center;line-height:1.6">
      Sent by <a href="https://eggypdf.com" style="color:#f5a623;text-decoration:none">EggyPDF</a>
    </p>
  </div>
</body></html>"""

    brevo_payload = {
        "sender": {"name": "EggyPDF", "email": SENDER_EMAIL},
        "to": [{"email": email, "name": name}],
        "replyTo": {"email": SENDER_EMAIL, "name": "EggyPDF"},
        "subject": f"Your CV is ready — {name}",
        "htmlContent": email_body
    }

    if pdf_b64:
        brevo_payload["attachment"] = [{
            "name": f"{safe_name}_Resume.pdf",
            "content": pdf_b64
        }]

    try:
        payload = json_lib.dumps(brevo_payload).encode("utf-8")
        req = urllib.request.Request(
            "https://api.brevo.com/v3/smtp/email",
            data=payload,
            headers={"Content-Type": "application/json", "api-key": BREVO_API_KEY},
            method="POST"
        )
        with urllib.request.urlopen(req, timeout=20) as resp:
            result = json_lib.loads(resp.read().decode("utf-8"))
            msg_id = result.get("messageId", "")
            print(f"Email sent OK: {msg_id}, pdf={'YES' if pdf_b64 else 'NO'}")
            return jsonify({"success": True, "pdf_attached": pdf_b64 is not None})
    except urllib.error.HTTPError as e:
        body = e.read().decode("utf-8", errors="ignore")
        print(f"Brevo error {e.code}: {body}")
        return jsonify({"error": f"Email failed: {e.code}"}), 500
    except Exception as e:
        return jsonify({"error": f"Email failed: {str(e)}"}), 500


# ─── AI: DIAGNOSTIC — LIST AVAILABLE MODELS ───
@app.route('/api/ai-models', methods=['GET', 'OPTIONS'])
def list_ai_models():
    """Lists all Gemini models available for the configured API key."""
    import urllib.request
    import urllib.error
    import json as json_lib

    GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY', '').strip()
    if not GEMINI_API_KEY:
        return jsonify({"error": "GEMINI_API_KEY not set on Render"}), 503

    try:
        url = f"https://generativelanguage.googleapis.com/v1beta/models?key={GEMINI_API_KEY}"
        req = urllib.request.Request(url, headers={"Content-Type": "application/json"})
        with urllib.request.urlopen(req, timeout=15) as resp:
            data = json_lib.loads(resp.read().decode("utf-8"))
        models = [m["name"] for m in data.get("models", [])]
        return jsonify({"available_models": models, "count": len(models)})
    except urllib.error.HTTPError as e:
        body = e.read().decode("utf-8", errors="ignore")
        return jsonify({"error": f"HTTP {e.code}", "detail": body[:300]}), e.code
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# ─── AI: RESUME SUGGESTIONS ───
@app.route('/api/ai-suggestions', methods=['POST', 'OPTIONS'])
def ai_suggestions():
    GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY', '').strip()
    if not GEMINI_API_KEY:
        return jsonify({"error": "AI service is not configured. Please add GEMINI_API_KEY on Render."}), 503

    data = request.get_json(silent=True) or {}
    job_title    = (data.get('job_title') or '').strip()
    suggest_type = (data.get('type') or 'bullets').strip()

    if not job_title:
        return jsonify({"error": "Please provide a job title."}), 400

    # Build prompt
    if suggest_type == 'bullets':
        prompt = (
            f"Write exactly 4 professional resume bullet points for a {job_title}. "
            "Each bullet point must: "
            "1. Start with a strong action verb (Developed, Led, Managed, Reduced, Built, etc). "
            "2. Be between 10 and 18 words long. "
            "3. Include a specific achievement, metric, or responsibility. "
            "4. Sound professional and ATS-friendly. "
            "Return ONLY a JSON array of exactly 4 strings. No extra text, no markdown, no explanation. "
            'Example: ["Developed scalable REST APIs serving 50,000 daily active users with 99.9 percent uptime", "Led a team of 6 engineers to deliver 3 product features ahead of schedule", "Reduced page load time by 45 percent through code optimization and lazy loading techniques", "Managed end-to-end deployment pipeline cutting release time from 2 days to 4 hours"]'
        )
    elif suggest_type == 'summary':
        prompt = (
            f'Write 4 professional resume summary paragraphs for a {job_title}. '
            'Each paragraph: 2 sentences, under 50 words, professional tone. '
            'Output format: JSON array of 4 strings only. '
            'No markdown. No explanation. No numbering. Just the JSON array.'
        )
    elif suggest_type == 'skills':
        prompt = (
            f"List the 8 most in-demand skills for a {job_title} in 2024. "
            "Return ONLY a valid JSON array of 8 short skill name strings. "
            'No explanation, no markdown. Example: ["Python", "SQL", "Excel"]'
        )
    else:
        return jsonify({"error": "Invalid type. Use bullets, summary, or skills."}), 400

    try:
        import urllib.request
        import urllib.error
        import json as json_lib

        payload_dict = {
            "contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {"temperature": 0.7, "maxOutputTokens": 2048}
        }

        # Using confirmed available models — gemini-2.5-flash confirmed working
        combos = [
            ("v1beta", "gemini-2.5-flash"),
            ("v1beta", "gemini-2.0-flash"),
            ("v1beta", "gemini-1.5-flash"),
            ("v1beta", "gemini-1.5-flash-latest"),
        ]

        result = None
        errors = []

        for api_ver, model in combos:
            try:
                url = (f"https://generativelanguage.googleapis.com"
                       f"/{api_ver}/models/{model}:generateContent"
                       f"?key={GEMINI_API_KEY}")
                payload = json_lib.dumps(payload_dict).encode("utf-8")
                req = urllib.request.Request(
                    url, data=payload,
                    headers={"Content-Type": "application/json"},
                    method="POST"
                )
                with urllib.request.urlopen(req, timeout=20) as resp:
                    result = json_lib.loads(resp.read().decode("utf-8"))
                break
            except urllib.error.HTTPError as e:
                if e.code == 400:
                    body = e.read().decode("utf-8", errors="ignore")
                    return jsonify({"error": f"API key rejected (400): {body[:200]}"}), 400
                if e.code == 429:
                    return jsonify({"error": "Daily AI limit reached. Try again tomorrow."}), 429
                errors.append(f"{api_ver}/{model}={e.code}")
                continue
            except Exception as ex:
                errors.append(f"{api_ver}/{model}={str(ex)[:40]}")
                continue

        if result is None:
            return jsonify({"error": f"All Gemini models failed: {'; '.join(errors)}"}), 500

        # Safely get text from response
        candidates = result.get("candidates", [])
        if not candidates:
            raise ValueError("Gemini returned no candidates")

        content_parts = candidates[0].get("content", {}).get("parts", [])
        if not content_parts:
            raise ValueError("Gemini returned empty content")

        text = content_parts[0].get("text", "").strip()
        if not text:
            raise ValueError("Gemini returned empty text")

        # Clean markdown fences
        text = text.replace("```json", "").replace("```", "").strip()

        # Strategy 1: direct JSON parse
        suggestions = None
        try:
            suggestions = json_lib.loads(text)
        except Exception:
            pass

        # Strategy 2: extract JSON array with regex
        if not suggestions:
            import re as _re2
            m = _re2.search(r'\[.*?\]', text, _re2.DOTALL)
            if m:
                try:
                    suggestions = json_lib.loads(m.group(0))
                except Exception:
                    pass

        # Strategy 3: split by newlines and clean each line
        if not suggestions:
            import re as _re3
            lines = []
            for line in text.splitlines():
                line = line.strip()
                # Remove leading numbers, bullets, quotes
                line = _re3.sub(r'^[0-9. \-\*]+', '', line)
                line = line.strip().strip('"').strip("'").strip(',')
                if len(line) > 10:
                    lines.append(line)
            suggestions = lines

        if not suggestions:
            raise ValueError("Could not parse AI response into suggestions")

        # Clean and validate each suggestion
        suggestions = [str(s).strip() for s in suggestions if s and len(str(s).strip()) > 8][:8]

        if not suggestions:
            raise ValueError("All suggestions were empty after cleaning")

        return jsonify({"suggestions": suggestions})

    except Exception as e:
        return jsonify({"error": f"AI failed: {str(e)}"}), 500


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 10000))
    app.run(host='0.0.0.0', port=port, debug=False)

